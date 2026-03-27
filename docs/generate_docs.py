#!/usr/bin/env python3
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Pt
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfgen import canvas
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer


ROOT = Path(__file__).resolve().parent
MANUAL_MD = ROOT / "GestureSlides_User_Manual.md"
MANUAL_PDF = ROOT / "GestureSlides_User_Manual.pdf"
PPTX_FILE = ROOT / "GestureSlides_Project_Presentation.pptx"
PRESENTATION_PDF = ROOT / "GestureSlides_Project_Presentation.pdf"


SLIDES = [
    {
        "title": "Presentation with gesture",
        "subtitle": "Advanced Presentation Studio with Multimodal Controls",
        "bullets": [
            "Elite editor + advanced design modules",
            "Presentation player + stop/zoom + two-finger/clap control",
            "Built with React, Fabric.js, Express, MediaPipe, OpenCV",
        ],
    },
    {
        "title": "Problem Statement",
        "bullets": [
            "Traditional presentations rely on keyboard or remote clickers",
            "Design-heavy tools are powerful but not touchless",
            "Need a demo app combining design flexibility with gesture control",
        ],
    },
    {
        "title": "Project Objectives",
        "bullets": [
            "Deliver richer Canva-style design workflow",
            "Add robust gesture set beyond basic navigation",
            "Support sound-triggered slide progression",
            "Provide stable save/load and polished demo flow",
        ],
    },
    {
        "title": "System Architecture",
        "bullets": [
            "Frontend: React + Fabric canvas editor and presentation overlay",
            "Backend: Node/Express APIs for save/load JSON",
            "Gesture Module: Python + MediaPipe + OpenCV",
            "Realtime link: WebSocket events to presentation controller",
        ],
    },
    {
        "title": "Core Modules",
        "bullets": [
            "Design modules: templates, themes, shapes, styling, layers",
            "Productivity modules: undo/redo, group/ungroup, duplicate, align",
            "Presentation mode: fullscreen, pause mode, zoom controls",
            "Gesture module: swipe, start, stop, two-finger zoom permission, thumbs-up zoom-out, clap control",
        ],
    },
    {
        "title": "Technology Stack",
        "bullets": [
            "Frontend: React, Vite, Fabric.js, CSS",
            "Backend: Node.js, Express",
            "Computer Vision: MediaPipe Tasks, OpenCV",
            "Protocol: WebSocket for gesture events",
        ],
    },
    {
        "title": "Data Model",
        "bullets": [
            "Presentation -> slides[] -> objects[]",
            "Each object stores type and Fabric properties",
            "Decks saved as JSON via /api/presentation/:id",
            "Supports local export for backup and portability",
        ],
    },
    {
        "title": "Gesture Pipeline",
        "bullets": [
            "Webcam frame capture",
            "Hand landmark detection (21 points)",
            "Smoothed motion + velocity/displacement thresholds",
            "Emit events: next / prev / start / stop / zoom_arm / zoom_in / zoom_out",
            "Audio detector: one clap next, two quick claps previous",
        ],
    },
    {
        "title": "Gesture Training Module",
        "bullets": [
            "Interactive trainer: python train_gesture_profile.py",
            "Calibrates personal swipe/zoom/sound thresholds",
            "Maps each primitive gesture to an operation",
            "Saves reusable gesture_profile.json",
        ],
    },
    {
        "title": "Noise Reduction Improvements",
        "bullets": [
            "Moving average smoothing window",
            "Swipe direction confirmation streak",
            "Cooldown + palm open-close sequence logic for start",
            "Pinch velocity gating for zoom stability",
            "Adaptive sound-floor threshold for clap events",
        ],
    },
    {
        "title": "How To Run (3 Terminals)",
        "bullets": [
            "Backend: cd backend && npm run dev",
            "Frontend: cd frontend && npm run dev",
            "Gesture: cd gesture && source .venv/bin/activate && python gesture_controller.py",
            "Open app at http://localhost:5173",
        ],
    },
    {
        "title": "Live Demo Flow",
        "bullets": [
            "Create 3 slides in editor",
            "Save and reload deck",
            "Start presentation mode",
            "Navigate with swipe + two-finger permission + pinch zoom + thumbs-up zoom-out + clap next/prev",
        ],
    },
    {
        "title": "Challenges and Fixes",
        "bullets": [
            "MediaPipe API differences (solutions vs tasks)",
            "WebSocket port conflict resolution",
            "OpenCV preview failures handled with headless fallback",
            "Gesture false positives reduced with stricter logic",
        ],
    },
    {
        "title": "Future Enhancements",
        "bullets": [
            "Realtime collaboration and shared cursors",
            "AI slide drafting + layout suggestions",
            "Gesture personalization profile",
            "Native PPTX animation timeline support",
        ],
    },
    {
        "title": "Conclusion",
        "bullets": [
            "Presentation with gesture demonstrates design + CV + audio integration",
            "Prototype is stable enough for project demo",
            "Architecture is extensible toward production features",
            "Thank you",
        ],
    },
]


def generate_pptx() -> None:
    prs = Presentation()
    for i, slide_data in enumerate(SLIDES):
        layout = prs.slide_layouts[0] if i == 0 else prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = slide_data["title"]

        if i == 0:
            subtitle = slide.placeholders[1]
            subtitle.text = slide_data["subtitle"]
            text_frame = subtitle.text_frame
            p = text_frame.add_paragraph()
            p.text = ""
            for bullet in slide_data["bullets"]:
                para = text_frame.add_paragraph()
                para.text = f"• {bullet}"
                para.level = 0
        else:
            body = slide.shapes.placeholders[1]
            text_frame = body.text_frame
            text_frame.clear()
            for idx, bullet in enumerate(slide_data["bullets"]):
                para = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
                para.text = bullet
                para.level = 0
                para.font.size = Pt(24 if idx == 0 else 21)

    prs.save(PPTX_FILE)


def _manual_story():
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "ManualTitle",
        parent=styles["Title"],
        fontSize=24,
        leading=30,
        spaceAfter=16,
    )
    h2 = ParagraphStyle(
        "ManualHeading",
        parent=styles["Heading2"],
        fontSize=14,
        leading=18,
        spaceBefore=10,
        spaceAfter=6,
    )
    body = ParagraphStyle(
        "ManualBody",
        parent=styles["BodyText"],
        fontSize=10.5,
        leading=14,
        spaceAfter=6,
    )
    bullet = ParagraphStyle(
        "ManualBullet",
        parent=styles["BodyText"],
        fontSize=10.5,
        leading=14,
        leftIndent=14,
        bulletIndent=6,
        spaceAfter=4,
    )

    story = [Paragraph("Presentation with gesture User Manual", title_style)]
    lines = MANUAL_MD.read_text(encoding="utf-8").splitlines()

    for raw in lines:
        line = raw.strip()
        if not line:
            story.append(Spacer(1, 0.12 * cm))
            continue
        if line.startswith("# "):
            continue
        if line.startswith("## "):
            story.append(Paragraph(line[3:], h2))
            continue
        if line.startswith("- "):
            story.append(Paragraph(line[2:], bullet, bulletText="•"))
            continue
        if line.startswith("```"):
            continue
        story.append(Paragraph(line, body))
    return story


def generate_manual_pdf() -> None:
    doc = SimpleDocTemplate(
        str(MANUAL_PDF),
        pagesize=A4,
        leftMargin=1.8 * cm,
        rightMargin=1.8 * cm,
        topMargin=1.8 * cm,
        bottomMargin=1.8 * cm,
    )
    doc.build(_manual_story())


def _draw_slide_page(c: canvas.Canvas, title: str, bullets: list[str], page_index: int, total_pages: int) -> None:
    width, height = landscape(A4)
    c.setFillColorRGB(0.97, 0.98, 1.0)
    c.rect(0, 0, width, height, stroke=0, fill=1)

    c.setFillColorRGB(0.1, 0.16, 0.25)
    c.setFont("Helvetica-Bold", 30)
    c.drawString(2 * cm, height - 2.4 * cm, title)

    y = height - 4.2 * cm
    c.setFont("Helvetica", 16)
    for bullet in bullets:
        text_obj = c.beginText(2.5 * cm, y)
        text_obj.textLine(f"• {bullet}")
        c.drawText(text_obj)
        y -= 1.15 * cm
        if y < 2.5 * cm:
            break

    c.setFillColorRGB(0.35, 0.4, 0.5)
    c.setFont("Helvetica", 10)
    c.drawRightString(width - 1.5 * cm, 1.2 * cm, f"Slide {page_index}/{total_pages}")


def generate_presentation_pdf() -> None:
    pdfmetrics.getRegisteredFontNames()
    c = canvas.Canvas(str(PRESENTATION_PDF), pagesize=landscape(A4))
    total = len(SLIDES)
    for idx, slide in enumerate(SLIDES, start=1):
        bullets = slide.get("bullets", [])
        if idx == 1:
            bullets = [slide.get("subtitle", "")] + bullets
        _draw_slide_page(c, slide["title"], bullets, idx, total)
        c.showPage()
    c.save()


def main() -> None:
    if not MANUAL_MD.exists():
        raise FileNotFoundError(f"Manual source not found: {MANUAL_MD}")

    generate_pptx()
    generate_manual_pdf()
    generate_presentation_pdf()

    print("Generated files:")
    print(f"- {PPTX_FILE}")
    print(f"- {MANUAL_PDF}")
    print(f"- {PRESENTATION_PDF}")


if __name__ == "__main__":
    main()
